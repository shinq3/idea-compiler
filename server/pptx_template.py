#!/usr/bin/env python3
"""
PPTX Template Generator - Rich Shape Style
stdin: JSON with cover + slides data
stdout: base64-encoded PPTX
"""
import sys
import json
import copy
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import os

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "../documents/frontier_template1.pptx")

# カラーパレット
C_PRIMARY   = RGBColor(0x66, 0x7e, 0xea)
C_SECONDARY = RGBColor(0x76, 0x4b, 0xa2)
C_SUCCESS   = RGBColor(0x48, 0xbb, 0x78)
C_WARNING   = RGBColor(0xed, 0x89, 0x36)
C_DANGER    = RGBColor(0xe5, 0x3e, 0x3e)
C_ACCENT    = RGBColor(0xe8, 0x3a, 0x3a)
C_TEXT      = RGBColor(0x2d, 0x37, 0x48)
C_MUTED     = RGBColor(0x71, 0x80, 0x96)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_BG   = RGBColor(0xF7, 0xFA, 0xFC)
C_CARD_BG2  = RGBColor(0xED, 0xF2, 0xF7)

PALETTE = [C_PRIMARY, C_SECONDARY, C_SUCCESS, C_WARNING, C_DANGER]

def replace_placeholder(slide, placeholder, value):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if placeholder not in para.text:
                continue
            if not para.runs:
                continue
            full_text = para.text.replace(placeholder, value)
            para.runs[0].text = full_text
            for run in para.runs[1:]:
                run.text = ""

def duplicate_slide(prs, slide_index):
    template = prs.slides[slide_index]
    slide = prs.slides.add_slide(template.slide_layout)
    sp_tree = slide.shapes._spTree
    for el in list(sp_tree):
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            sp_tree.remove(el)
    for shape in template.shapes:
        el = copy.deepcopy(shape.element)
        if shape.shape_type == 13:
            try:
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                r_embed_key = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for blip in el.iter('{%s}blip' % ns):
                    old_rId = blip.get(r_embed_key)
                    if old_rId:
                        img_part = template.part.related_part(old_rId)
                        new_rId = slide.part.relate_to(img_part, RT.IMAGE)
                        blip.set(r_embed_key, new_rId)
            except Exception:
                pass
        sp_tree.append(el)
    return slide

# ===== 図形描画ヘルパー =====

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=14, bold=False,
             color=None, align='left', valign='top', wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    if align == 'center':
        from pptx.enum.text import PP_ALIGN
        para.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        from pptx.enum.text import PP_ALIGN
        para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or C_TEXT
    return txBox

def add_multiline_text(slide, items, x, y, w, h, font_size=13, color=None, bullet_color=None, space_before=5):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_before = Pt(space_before)
        dot = para.add_run()
        dot.text = "● "
        dot.font.size = Pt(9)
        dot.font.color.rgb = bullet_color or C_PRIMARY
        txt = para.add_run()
        txt.text = item
        txt.font.size = Pt(font_size)
        txt.font.color.rgb = color or C_TEXT

# ===== レイアウト別コンテンツ描画 =====

CX, CY, CW, CH = 0.75, 1.9, 11.83, 4.8

def layout_bullets(slide, data):
    """左アクセントバー付き大カード + カラー箇条書き"""
    bullets = data.get("bullets", [])
    # メインカード
    add_rounded_rect(slide, CX, CY, CW, CH, C_CARD_BG,
                     line_color=RGBColor(0xE2, 0xE8, 0xF0))
    # 左アクセントバー
    add_rect(slide, CX, CY + 0.1, 0.07, CH - 0.2, C_PRIMARY)
    # 箇条書き
    add_multiline_text(slide, bullets, CX + 0.25, CY + 0.2, CW - 0.4, CH - 0.4,
                       font_size=15, bullet_color=C_PRIMARY)

def layout_two_column(slide, data):
    """2カラムカード（左青・右紫）"""
    columns = data.get("columns", [])
    col_w = 5.6
    gap = 0.43
    colors = [C_PRIMARY, C_SECONDARY]
    for ci, col in enumerate(columns[:2]):
        cx = CX + (col_w + gap) * ci
        color = colors[ci]
        # カード背景
        add_rounded_rect(slide, cx, CY, col_w, CH, C_CARD_BG,
                         line_color=RGBColor(0xE2, 0xE8, 0xF0))
        # ヘッダー帯
        add_rounded_rect(slide, cx, CY, col_w, 0.55, color)
        # 見出しテキスト
        add_text(slide, col.get("heading", ""), cx + 0.15, CY + 0.1,
                 col_w - 0.3, 0.4, font_size=14, bold=True, color=C_WHITE, align='center')
        # 箇条書き
        add_multiline_text(slide, col.get("points", []),
                           cx + 0.15, CY + 0.7, col_w - 0.3, CH - 0.85,
                           font_size=14, bullet_color=color)

def layout_metrics(slide, data):
    """数値メトリクスカード（3列）"""
    metrics = data.get("metrics", [])[:3]
    n = len(metrics)
    if n == 0:
        return
    card_w = (CW - (n - 1) * 0.3) / n
    colors = [C_PRIMARY, C_SUCCESS, C_WARNING]
    for mi, m in enumerate(metrics):
        cx = CX + (card_w + 0.3) * mi
        color = colors[mi % len(colors)]
        # カード
        add_rounded_rect(slide, cx, CY, card_w, CH, C_CARD_BG,
                         line_color=RGBColor(0xE2, 0xE8, 0xF0))
        # 上部カラーバー
        add_rounded_rect(slide, cx, CY, card_w, 0.08, color)
        # アイコン/絵文字
        icon = m.get("icon", "📊")
        add_text(slide, icon, cx, CY + 0.25, card_w, 0.9,
                 font_size=36, align='center')
        # 数値
        add_text(slide, m.get("value", ""), cx, CY + 1.2, card_w, 0.9,
                 font_size=32, bold=True, color=color, align='center')
        # ラベル
        add_text(slide, m.get("label", ""), cx, CY + 2.1, card_w, 0.5,
                 font_size=13, bold=True, color=C_TEXT, align='center')
        # 説明
        desc = m.get("description", "")
        if desc:
            add_text(slide, desc, cx + 0.2, CY + 2.7, card_w - 0.4, 1.8,
                     font_size=11, color=C_MUTED, align='center', wrap=True)

def layout_table(slide, data):
    """スタイル付きテーブル"""
    rows_data = data.get("rows", [])
    if not rows_data:
        return
    cols = len(rows_data[0])
    n_rows = len(rows_data)
    tbl_h = min(n_rows * 0.5, CH)
    table = slide.shapes.add_table(
        n_rows, cols,
        Inches(CX), Inches(CY),
        Inches(CW), Inches(tbl_h)
    ).table
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            tf = cell.text_frame
            para = tf.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(12)
            run.font.bold = (ri == 0)
            if ri == 0:
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_CARD_BG
                run.font.color.rgb = C_TEXT
            else:
                run.font.color.rgb = C_TEXT

def layout_steps(slide, data):
    """ステップ/タイムライン（番号付きカード）"""
    steps = data.get("steps", [])[:4]
    n = len(steps)
    if n == 0:
        return
    card_w = (CW - (n - 1) * 0.25) / n
    colors = [C_PRIMARY, C_SECONDARY, C_SUCCESS, C_WARNING]
    for si, step in enumerate(steps):
        cx = CX + (card_w + 0.25) * si
        color = colors[si % len(colors)]
        # カード背景
        add_rounded_rect(slide, cx, CY, card_w, CH, C_CARD_BG,
                         line_color=RGBColor(0xE2, 0xE8, 0xF0))
        # 番号円（小さい色付き矩形で代用）
        add_rounded_rect(slide, cx + card_w/2 - 0.3, CY + 0.2, 0.6, 0.6, color)
        add_text(slide, str(si + 1), cx + card_w/2 - 0.3, CY + 0.22,
                 0.6, 0.5, font_size=18, bold=True, color=C_WHITE, align='center')
        # タイトル
        add_text(slide, step.get("title", ""), cx + 0.1, CY + 0.95,
                 card_w - 0.2, 0.5, font_size=13, bold=True, color=color, align='center')
        # 説明
        add_text(slide, step.get("description", ""), cx + 0.15, CY + 1.55,
                 card_w - 0.3, CH - 1.8, font_size=12, color=C_TEXT, wrap=True)

def add_content_to_slide(slide, content_type, data):
    if content_type == "bullets":
        layout_bullets(slide, data)
    elif content_type == "two_column":
        layout_two_column(slide, data)
    elif content_type == "metrics":
        layout_metrics(slide, data)
    elif content_type == "table":
        layout_table(slide, data)
    elif content_type == "steps":
        layout_steps(slide, data)
    else:
        layout_bullets(slide, data)

def update_slide_number(slide, number):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip().isdigit():
                    run.text = str(number)
                    return

def generate(data):
    prs = Presentation(TEMPLATE_PATH)
    cover = data.get("cover", {})
    slides_data = data.get("slides", [])

    # スライド1: 表紙
    cover_slide = prs.slides[0]
    replace_placeholder(cover_slide, "{Title}", cover.get("title", ""))
    replace_placeholder(cover_slide, "{Message}", cover.get("message", ""))
    client_str = f"{cover.get('client', '')} / {cover.get('project_name', '')}"
    replace_placeholder(cover_slide, "{Client} / {Project Name}", client_str)
    replace_placeholder(cover_slide, "{Client}", cover.get("client", ""))
    replace_placeholder(cover_slide, "{Project Name}", cover.get("project_name", ""))
    replace_placeholder(cover_slide, "{date}", cover.get("date", ""))

    # スライド2以降: 内容ページ複製
    for i, slide_data in enumerate(slides_data):
        new_slide = duplicate_slide(prs, 1)
        replace_placeholder(new_slide, "{Key message}", slide_data.get("key_message", ""))
        replace_placeholder(new_slide, "{Page Title}", slide_data.get("page_title", ""))
        add_content_to_slide(new_slide, slide_data.get("content_type", "bullets"), slide_data)
        update_slide_number(new_slide, i + 2)

    # テンプレートのスライド2（空のマスター）を削除
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[1])

    output = io.BytesIO()
    prs.save(output)
    return base64.b64encode(output.getvalue()).decode("utf-8")

if __name__ == "__main__":
    raw = sys.stdin.read()
    data = json.loads(raw)
    result = generate(data)
    print(result)
